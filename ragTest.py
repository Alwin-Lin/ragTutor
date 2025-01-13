from dotenv import load_dotenv
import streamlit as st
from google.generativeai import GenerativeModel, configure
from sentence_transformers import SentenceTransformer
import numpy as np
from typing import List, Dict, Tuple, Union
import pandas as pd
from sklearn.metrics.pairwise import cosine_similarity
import os
import pickle
from pathlib import Path
import json
from pptx import Presentation
from datetime import datetime
import google.generativeai as genai
import PyPDF2
from docx import Document
from dataclasses import dataclass
from tqdm import tqdm

# Load environment variables
load_dotenv()

# Constants
COURSE_FOLDER = "./ENG3004_222"
VECTOR_DB_PATH = "vector_db"

@dataclass
class ProcessedDocument:
    """Data class to store processed document information"""
    content: Union[str, List[str]]
    source_path: str
    doc_type: str
    quality_label: str = None
    assignment_number: str = None

class DocumentProcessor:
    def read_pdf(self, file_path: str) -> List[Dict[str, any]]:
        """Extract text from PDF files with page tracking"""
        pages = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    try:
                        text = page.extract_text()
                        if text.strip():
                            pages.append({
                                'content': text,
                                'page_number': page_num
                            })
                    except Exception as e:
                        self.error_logs.append(f"Error in {file_path}, page {page_num}: {str(e)}")
                        continue
        except Exception as e:
            self.error_logs.append(f"Error processing {file_path}: {str(e)}")
        return pages

    def read_pptx(self, file_path: str) -> List[Dict[str, any]]:
        """Extract text from PowerPoint files with slide tracking"""
        slides = []
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                if slide_text:
                    slides.append({
                        'content': "\n".join(slide_text),
                        'slide_number': slide_num
                    })
        except Exception as e:
            self.error_logs.append(f"Error processing {file_path}: {str(e)}")
        return slides

    def read_docx(self, file_path: str) -> List[Dict[str, any]]:
        """Extract text from DOCX files with paragraph tracking"""
        sections = []
        try:
            doc = Document(file_path)
            current_section = []
            section_count = 1
            
            for para in doc.paragraphs:
                if para.text.strip():
                    current_section.append(para.text)
                    if len("\n".join(current_section)) >= 1000:
                        sections.append({
                            'content': "\n".join(current_section),
                            'section': section_count
                        })
                        current_section = []
                        section_count += 1
            
            if current_section:
                sections.append({
                    'content': "\n".join(current_section),
                    'section': section_count
                })
        except Exception as e:
            self.error_logs.append(f"Error processing {file_path}: {str(e)}")
        return sections

    def chunk_text(self, content: Dict[str, any], chunk_size: int = 1000) -> List[Dict[str, any]]:
        """Split text into smaller chunks while preserving source information"""
        if not content['content']:
            return []
            
        chunks = []
        words = content['content'].split()
        current_chunk = []
        current_size = 0
        
        for word in words:
            if current_size + len(word) > chunk_size:
                if current_chunk:
                    chunks.append({
                        'content': " ".join(current_chunk),
                        **{k: v for k, v in content.items() if k != 'content'}
                    })
                current_chunk = [word]
                current_size = len(word)
            else:
                current_chunk.append(word)
                current_size += len(word) + 1
        
        if current_chunk:
            chunks.append({
                'content': " ".join(current_chunk),
                **{k: v for k, v in content.items() if k != 'content'}
            })
        
        return chunks

class VectorStore:
    def __init__(self, base_path: str):
        self.base_path = Path(base_path)
        self.base_path.mkdir(exist_ok=True, parents=True)
        
    def save_data(self, 
                  embeddings: Dict,
                  contents: Dict,
                  metadata: Dict,
                  timestamp: str):
        """Save all vector store data"""
        data = {
            'embeddings': embeddings,
            'contents': contents,
            'metadata': metadata,
            'timestamp': timestamp
        }
        
        temp_file = self.base_path / "vector_store.tmp"
        final_file = self.base_path / "vector_store.pkl"
        
        # Save to temporary file first
        with open(temp_file, "wb") as f:
            pickle.dump(data, f)
        
        # Rename to final file (atomic operation)
        temp_file.replace(final_file)
    
    def load_data(self) -> Tuple[Dict, Dict, Dict, str]:
        """Load vector store data"""
        try:
            with open(self.base_path / "vector_store.pkl", "rb") as f:
                data = pickle.load(f)
                return (data.get('embeddings', {}), 
                       data.get('contents', {}), 
                       data.get('metadata', {}), 
                       data.get('timestamp', ""))
        except (FileNotFoundError, pickle.UnpicklingError, EOFError) as e:
            print(f"Error loading vector store: {str(e)}")
            return {}, {}, {}, ""

class RAGTutor:
    def __init__(self, api_key: str):
        # Initialize vector store
        self.vector_store = VectorStore(VECTOR_DB_PATH)
        self.processor = DocumentProcessor()
        
        # Configure Gemini API with provided key
        if not api_key:
            raise ValueError("Google API Key is required")
        configure(api_key=api_key)
        
        # Initialize models
        try:
            self.embed_model = SentenceTransformer('all-MiniLM-L6-v2')
            self.llm = GenerativeModel('gemini-pro')
        except Exception as e:
            raise RuntimeError(f"Error initializing models: {str(e)}")
        
        # Initialize storage
        self.document_embeddings = {}
        self.document_contents = {}
        self.document_metadata = {}
        self.last_updated = ""
        
        # Load or create vector database
        self.load_or_create_vectors()

    def load_or_create_vectors(self):
        """Load existing vector database or create new one"""
        (self.document_embeddings, 
         self.document_contents,
         self.document_metadata,
         self.last_updated) = self.vector_store.load_data()
        
        # If no existing database or it's too old, recreate it
        if not self.document_embeddings or self.should_update_vectors():
            self.process_all_documents()
            self.save_vectors()

    def should_update_vectors(self) -> bool:
        """Check if vectors should be updated based on file modifications"""
        if not self.last_updated:
            return True
            
        last_update = datetime.strptime(self.last_updated, "%Y-%m-%d")
        course_folder = Path(COURSE_FOLDER)
        
        # Check if any course files were modified after last update
        for path in course_folder.rglob("*"):
            if path.is_file():
                mtime = datetime.fromtimestamp(path.stat().st_mtime)
                if mtime > last_update:
                    return True
        
        return False
    
    def save_vectors(self):
        """Save current vector database"""
        self.vector_store.save_data(
            self.document_embeddings,
            self.document_contents,
            self.document_metadata,
            datetime.now().strftime("%Y-%m-%d")
        )

    def process_all_documents(self):
        """Process all documents in the folder structure"""
        course_folder = Path(COURSE_FOLDER)
        if not course_folder.exists():
            raise FileNotFoundError(f"Course folder not found: {COURSE_FOLDER}")
        
        # Process lecture notes
        self._process_folder(course_folder / "Lecture notes", "lecture")
        
        # Process tutorial sheets
        self._process_folder(course_folder / "Tutorial sheets", "tutorial")
        
        # Process assignments
        for i in range(1, 3):  # Assignments 1 and 2
            assignment_folder = course_folder / f"Assignment{i}"
            if assignment_folder.exists():
                self._process_assignment_folder(assignment_folder, f"Assignment{i}")
        
        # Process rubric
        rubric_path = course_folder / "rubric.docx"
        if rubric_path.exists():
            self._process_single_file(rubric_path, "rubric")
    
    def _process_assignment_folder(self, folder: Path, assignment_number: str):
        """Process assignment folder including examples"""
        # Process assignment prompt
        prompt_file = folder / f"{assignment_number.lower()}.pdf"
        if prompt_file.exists():
            self._process_single_file(prompt_file, "assignment_prompt", 
                                   assignment_number=assignment_number)
        
        # Process examples
        for quality in ['Poor', 'Average', 'Good']:
            quality_folder = folder / quality
            if quality_folder.exists():
                for example_file in quality_folder.glob("*.pdf"):
                    self._process_single_file(
                        example_file,
                        "student_example",
                        quality_label=quality.lower(),
                        assignment_number=assignment_number
                    )

    def _process_folder(self, folder_path: Path, category: str):
        """Process all files in a folder with progress bar"""
        if not folder_path.exists():
            st.warning(f"Folder {folder_path} does not exist")
            return
        
        files = list(folder_path.glob("*"))
        for file_path in tqdm(files, desc=f"Processing {category} files"):
            self._process_single_file(file_path, category)                
    
    def _process_single_file(self, file_path: Path, doc_type: str,
                        quality_label: str = None,
                        assignment_number: str = None):
        """Process a single file with source tracking"""
        try:
            content_sections = []
            if file_path.suffix.lower() == '.pdf':
                content_sections = self.processor.read_pdf(str(file_path))
                print(file_path.suffix.lower())
            elif file_path.suffix.lower() == '.docx':
                content_sections = self.processor.read_docx(str(file_path))
                print(file_path.suffix.lower())
            elif file_path.suffix.lower() in ['.pptx', '.ppt']:
                content_sections = self.processor.read_pptx(str(file_path))
                print(file_path.suffix.lower())
            
            if content_sections:
                doc_id = self._generate_doc_id(file_path, doc_type, 
                                            assignment_number, quality_label)
                
                all_chunks = []
                for section in content_sections:
                    # Ensure section is a dictionary
                    if isinstance(section, str):
                        section = {'content': section}
                    
                    # Add metadata to section
                    section['file_name'] = file_path.name
                    section['doc_type'] = doc_type
                    section['quality_label'] = quality_label
                    section['assignment_number'] = assignment_number
                    
                    chunks = self.processor.chunk_text(section)
                    all_chunks.extend(chunks)
                
                if all_chunks:
                    self._add_to_index(doc_id, all_chunks)
        
        except Exception as e:
            self.processor.error_logs.append(f"Error processing {file_path}: {str(e)}")
            print(f"Error details: {str(e)}")
    
    #ToDo: What does this do, why is it here?
    ### doc_id is for preventing simular document name and for labeling the quality of student works
    def _generate_doc_id(self, file_path: Path, doc_type: str,
                        assignment_number: str = None,
                        quality_label: str = None) -> str:
        """Generate unique document ID"""
        doc_id = f"{doc_type}_{file_path.stem}"
        if assignment_number:
            doc_id = f"{assignment_number}_{doc_id}"
        if quality_label:
            doc_id = f"{doc_id}_{quality_label}"
        print(doc_id)
        return doc_id
    
    def _add_to_index(self, doc_id: str, chunks: List[Dict[str, any]]):
        """Add document to searchable index with source information"""
        try:
            # Extract just the text content for embedding
            texts = [chunk['content'] if isinstance(chunk, dict) else chunk for chunk in chunks]
            embeddings = self.embed_model.encode(texts)
            
            self.document_embeddings[doc_id] = embeddings
            self.document_contents[doc_id] = chunks
            
            # Store simplified metadata
            self.document_metadata[doc_id] = {
                'type': doc_id.split('_')[0],  # Extract type from doc_id
                'quality_label': None,
                'assignment_number': None
            }
            
            # Add additional metadata if available
            if '_Assignment' in doc_id:
                assignment_part = doc_id.split('_Assignment')[1]
                assignment_number = 'Assignment' + assignment_part[0]  # Get the number
                self.document_metadata[doc_id]['assignment_number'] = assignment_number
                
            # Add quality label if present
            if any(label in doc_id.lower() for label in ['good', 'average', 'poor']):
                for label in ['good', 'average', 'poor']:
                    if label in doc_id.lower():
                        self.document_metadata[doc_id]['quality_label'] = label
                        break
                        
        except Exception as e:
            self.processor.error_logs.append(f"Error embedding {doc_id}: {str(e)}")
    
    def retrieve_relevant_context(self, query: str, top_k: int = 3) -> List[Dict[str, any]]:
        """Retrieve most relevant content with explicit source information"""
        query_embedding = self.embed_model.encode(query)
        
        all_similarities = []
        for doc_id, embeddings in self.document_embeddings.items():
            similarities = cosine_similarity([query_embedding], embeddings)[0]
            for idx, sim in enumerate(similarities):
                all_similarities.append((sim, doc_id, idx))
        
        all_similarities.sort(reverse=True)
        relevant_content = []
        
        for sim, doc_id, idx in all_similarities[:top_k]:
            chunk = self.document_contents[doc_id][idx]
            metadata = self.document_metadata[doc_id]
            
            # Create detailed source information
            source_info = {
                'content': chunk['content'] if isinstance(chunk, dict) else chunk,
                'similarity_score': float(sim),
                'document_type': metadata['type'],
                'file_path': chunk.get('file_name', 'Unknown'),
                'location': {}
            }
            
            # Add specific location information based on document type
            if isinstance(chunk, dict):
                if 'page_number' in chunk:
                    source_info['location']['page'] = chunk['page_number']
                if 'slide_number' in chunk:
                    source_info['location']['slide'] = chunk['slide_number']
                if 'section' in chunk:
                    source_info['location']['section'] = chunk['section']
            
            # Add quality and assignment information if available
            if metadata.get('quality_label'):
                source_info['quality_label'] = metadata['quality_label']
            if metadata.get('assignment_number'):
                source_info['assignment_number'] = metadata['assignment_number']
            
            relevant_content.append(source_info)
        print(relevant_content)
        return relevant_content
    
    def format_source_citation(self, source_info: Dict) -> str:
        """Format source information into a readable citation"""
        citation = f"\n\nSource: {source_info['file_path']}"
        
        # Add location information
        if source_info['location']:
            locations = []
            if 'page' in source_info['location']:
                locations.append(f"Page {source_info['location']['page']}")
            if 'slide' in source_info['location']:
                locations.append(f"Slide {source_info['location']['slide']}")
            if 'section' in source_info['location']:
                locations.append(f"Section {source_info['location']['section']}")
            if locations:
                citation += f" ({', '.join(locations)})"
        
        # Add quality and assignment information for student examples
        if source_info.get('quality_label') and source_info.get('assignment_number'):
            citation += f"\nStudent Example - {source_info['assignment_number']}"
            citation += f"\nQuality Level: {source_info['quality_label'].title()}"
        
        return citation
    
    def generate_response(self, 
                         student_query: str,
                         student_work: str = None) -> Dict[str, any]:
        """Generate tutoring response with explicit source materials"""
        try:
            # Retrieve relevant context
            contexts = self.retrieve_relevant_context(
                student_work if student_work else student_query
            )
            
            # Prepare formatted context with citations
            formatted_contexts = []
            for ctx in contexts:
                formatted_context = {
                    'content': ctx['content'],
                    'citation': self.format_source_citation(ctx),
                    'relevance': ctx['similarity_score']
                }
                formatted_contexts.append(formatted_context)
            
            # Create prompt with explicit source references
            prompt_parts = ["You are an AI tutor for the course 'Society and the Engineer'.\n"]
            prompt_parts.append("\nRelevant course materials:")
            
            for i, ctx in enumerate(formatted_contexts, 1):
                prompt_parts.append(f"\n[Reference {i}]:")
                prompt_parts.append(ctx['content'])
                prompt_parts.append(ctx['citation'])
            
            if student_work:
                prompt_parts.append(f"\n\nStudent work submitted: {student_work}")
                
                if student_query == "Please provide feedback on this work and suggest improvements.":
                    prompt_parts.append("\n\nProvide detailed feedback that:")
                    prompt_parts.append("\n- Rates the work as poor/average/strong")
                    prompt_parts.append("\n- Analyzes strengths and areas for improvement")
                    prompt_parts.append("\n- References specific parts of the submission")
                    prompt_parts.append("\n- Suggests concrete improvements")
                    prompt_parts.append("\n- References specific course materials using [Reference X] format")
                else:
                    prompt_parts.append(f"\n\nStudent question: {student_query}")
            else:
                prompt_parts.append(f"\n\nStudent query: {student_query}")
            
            prompt_parts.append("\n\nIn your response:")
            prompt_parts.append("\n- Directly answer the question")
            prompt_parts.append("\n- Use [Reference X] to cite specific source materials")
            prompt_parts.append("\n- Connect ideas across different references")
            prompt_parts.append("\n- Provide specific examples when possible")
            
            # Generate response
            response = self.llm.generate_content("".join(prompt_parts))
            response_text = ""
            if hasattr(response, 'parts'):
                for part in response.parts:
                    response_text += part.text
            else:
                response_text = getattr(response, 'text', 'Unable to generate response')
            
            # Return both response and source materials
            return {
                'response': response_text,
                'sources': formatted_contexts
            }
            
        except Exception as e:
            print(f"Error generating response: {str(e)}")
            return {
                'response': "I apologize, but I encountered an error generating the response. Please try again.",
                'sources': []
            }

def create_tutor_ui():
    st.set_page_config(layout="wide", page_title="ENG3004 AI Tutor")
    st.title("Society and the Engineer - AI Tutor")
    
    # Add API key input at the top
    api_key = st.text_input(
        "Enter your Google API Key", 
        type="password",
        help="Get your API key from https://makersuite.google.com/app/apikey"
    )
    
    if not api_key:
        st.warning("Please enter your Google API Key to use the tutor.")
        st.markdown("""
        To get an API key:
        1. Go to [Google AI Studio](https://makersuite.google.com/app/apikey)
        2. Sign in or create an account
        3. Click 'Create API Key'
        4. Copy and paste your key above
        
        Your API key will only be used for this session and won't be stored.
        """)
        return
    
    # Initialize tutor with API key if not already initialized or if key changed
    if 'tutor' not in st.session_state or 'api_key' not in st.session_state or st.session_state.api_key != api_key:
        with st.spinner("Initializing tutor..."):
            try:
                st.session_state.tutor = RAGTutor(api_key)
                st.session_state.api_key = api_key
                st.success("Tutor initialized successfully!")
            except Exception as e:
                st.error(f"Error initializing tutor: {str(e)}")
                return
    
    # Input areas
    col1, col2 = st.columns([1, 1])
    
    with col1:
        student_query = st.text_area(
            "What's your question? (Optional if submitting work for feedback)", 
            height=150
        )
    
    with col2:
        student_work = st.text_area(
            "Paste your work here for feedback (Optional)", 
            height=150
        )
    
    if st.button("Get Help"):
        if not student_query and not student_work:
            st.warning("Please either ask a question or submit work for feedback.")
            return
            
        with st.spinner("Generating response..."):
            try:
                if not student_query and student_work:
                    default_query = "Please provide feedback on this work and suggest improvements."
                    result = st.session_state.tutor.generate_response(
                        default_query,
                        student_work
                    )
                else:
                    result = st.session_state.tutor.generate_response(
                        student_query,
                        student_work if student_work else None
                    )
                
                # Display main response
                st.write("### Response")
                st.write(result['response'])
                
                # Display source materials
                st.write("### Source Materials")
                for i, source in enumerate(result['sources'], 1):
                    with st.expander(f"Reference {i}"):
                        st.write(source['content'])
                        st.write(source['citation'])
                        st.write(f"Relevance Score: {source['relevance']:.2f}")
            except Exception as e:
                st.error(f"Error generating response: {str(e)}")
                st.warning("Please check if your API key is valid and has sufficient quota.")

if __name__ == "__main__":
    create_tutor_ui()