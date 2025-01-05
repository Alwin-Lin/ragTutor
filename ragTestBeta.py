from dotenv import load_dotenv
import streamlit as st
from google.generativeai import GenerativeModel, configure
from sentence_transformers import SentenceTransformer
import numpy as np
from typing import List, Dict, Tuple, Union
import pandas as pd
from sklearn.metrics.pairwise import cosine_similarity
import os
import pickle
from pathlib import Path
import json
from pptx import Presentation
from datetime import datetime
import google.generativeai as genai
import PyPDF2
from docx import Document
from dataclasses import dataclass
from tqdm import tqdm

# Load environment variables
load_dotenv()

# Constants
COURSE_FOLDER = "/Users/Alwin/Downloads/ENG3004_222"
VECTOR_DB_PATH = "vector_db"

@dataclass
class ProcessedDocument:
    """Data class to store processed document information"""
    content: Union[str, List[str]]
    source_path: str
    doc_type: str
    quality_label: str = None
    assignment_number: str = None

class DocumentProcessor:
    def __init__(self):
        self.error_logs = []
        
    def read_pdf(self, file_path: str) -> str:
        """Extract text from PDF files with error recovery"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        text += page.extract_text() + "\n"
                    except Exception as e:
                        self.error_logs.append(f"Error in {file_path}, page {page_num}: {str(e)}")
                        continue
        except Exception as e:
            self.error_logs.append(f"Error processing {file_path}: {str(e)}")
        return text

    def read_docx(self, file_path: str) -> str:
        """Extract text from DOCX files with error recovery"""
        try:
            doc = Document(file_path)
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            self.error_logs.append(f"Error processing {file_path}: {str(e)}")
            return ""

    def read_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint files with error recovery"""
        try:
            prs = Presentation(file_path)
            text = []
            
            for slide in prs.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                text.append("\n".join(slide_text))
            
            return "\n===SLIDE BREAK===\n".join(text)
        except Exception as e:
            self.error_logs.append(f"Error processing {file_path}: {str(e)}")
            return ""

    @staticmethod
    def chunk_text(text: str, chunk_size: int = 1000) -> List[str]:
        """Split text into smaller chunks while preserving slide breaks"""
        if not text:
            return []
            
        chunks = []
        slides = text.split("===SLIDE BREAK===")
        
        for slide in slides:
            words = slide.split()
            current_chunk = []
            current_size = 0
            
            for word in words:
                if current_size + len(word) > chunk_size:
                    if current_chunk:  # Only append if there's content
                        chunks.append(" ".join(current_chunk))
                    current_chunk = [word]
                    current_size = len(word)
                else:
                    current_chunk.append(word)
                    current_size += len(word) + 1
            
            if current_chunk:  # Only append if there's content
                chunks.append(" ".join(current_chunk))
        
        return [chunk for chunk in chunks if chunk.strip()]  # Remove empty chunks

class VectorStore:
    def __init__(self, base_path: str):
        self.base_path = Path(base_path)
        self.base_path.mkdir(exist_ok=True, parents=True)  # Added parents=True
        
    def save_data(self, 
                  embeddings: Dict,
                  contents: Dict,
                  metadata: Dict,
                  timestamp: str):
        """Save all vector store data"""
        data = {
            'embeddings': embeddings,
            'contents': contents,
            'metadata': metadata,
            'timestamp': timestamp
        }
        
        temp_file = self.base_path / "vector_store.tmp"
        final_file = self.base_path / "vector_store.pkl"
        
        # Save to temporary file first
        with open(temp_file, "wb") as f:
            pickle.dump(data, f)
        
        # Rename to final file (atomic operation)
        temp_file.replace(final_file)
    
    def load_data(self) -> Tuple[Dict, Dict, Dict, str]:
        """Load vector store data"""
        try:
            with open(self.base_path / "vector_store.pkl", "rb") as f:
                data = pickle.load(f)
                return (data.get('embeddings', {}), 
                       data.get('contents', {}), 
                       data.get('metadata', {}), 
                       data.get('timestamp', ""))
        except (FileNotFoundError, pickle.UnpicklingError, EOFError) as e:
            print(f"Error loading vector store: {str(e)}")
            return {}, {}, {}, ""

class RAGTutor:
    def __init__(self):
        # Initialize vector store
        self.vector_store = VectorStore(VECTOR_DB_PATH)
        self.processor = DocumentProcessor()
        
        # Configure Gemini API
        api_key = os.getenv('GOOGLE_API_KEY')
        if not api_key:
            raise ValueError("GOOGLE_API_KEY not found in environment variables")
        configure(api_key=api_key)
        
        # Initialize models
        try:
            self.embed_model = SentenceTransformer('all-MiniLM-L6-v2')
            self.llm = GenerativeModel('gemini-pro')
        except Exception as e:
            raise RuntimeError(f"Error initializing models: {str(e)}")
        
        # Initialize storage
        self.document_embeddings = {}
        self.document_contents = {}
        self.document_metadata = {}
        self.last_updated = ""
        
        # Load or create vector database
        self.load_or_create_vectors()

    def load_or_create_vectors(self):
        """Load existing vector database or create new one"""
        (self.document_embeddings, 
         self.document_contents,
         self.document_metadata,
         self.last_updated) = self.vector_store.load_data()
        
        # If no existing database or it's too old, recreate it
        if not self.document_embeddings or self.should_update_vectors():
            self.process_all_documents()
            self.save_vectors()

    def should_update_vectors(self) -> bool:
        """Check if vectors should be updated based on file modifications"""
        if not self.last_updated:
            return True
            
        last_update = datetime.strptime(self.last_updated, "%Y-%m-%d")
        course_folder = Path(COURSE_FOLDER)
        
        # Check if any course files were modified after last update
        for path in course_folder.rglob("*"):
            if path.is_file():
                mtime = datetime.fromtimestamp(path.stat().st_mtime)
                if mtime > last_update:
                    return True
        
        return False
    
    def save_vectors(self):
        """Save current vector database"""
        self.vector_store.save_data(
            self.document_embeddings,
            self.document_contents,
            self.document_metadata,
            datetime.now().strftime("%Y-%m-%d")
        )

    def process_all_documents(self):
        """Process all documents in the folder structure"""
        course_folder = Path(COURSE_FOLDER)
        if not course_folder.exists():
            raise FileNotFoundError(f"Course folder not found: {COURSE_FOLDER}")
        
        # Process lecture notes
        self._process_folder(course_folder / "Lecture Notes", "lecture")
        
        # Process tutorial sheets
        self._process_folder(course_folder / "Tutorial sheets", "tutorial")
        
        # Process assignments
        for i in range(1, 3):  # Assignments 1 and 2
            assignment_folder = course_folder / f"Assignment{i}"
            if assignment_folder.exists():
                self._process_assignment_folder(assignment_folder, f"Assignment{i}")
        
        # Process rubric
        rubric_path = course_folder / "rubric.docx"
        if rubric_path.exists():
            self._process_single_file(rubric_path, "rubric")
    
    def _process_assignment_folder(self, folder: Path, assignment_number: str):
        """Process assignment folder including examples"""
        # Process assignment prompt
        prompt_file = folder / f"{assignment_number.lower()}.pdf"
        if prompt_file.exists():
            self._process_single_file(prompt_file, "assignment_prompt", 
                                   assignment_number=assignment_number)
        
        # Process examples
        for quality in ['Poor', 'Average', 'Good']:
            quality_folder = folder / quality
            if quality_folder.exists():
                for example_file in quality_folder.glob("*.pdf"):
                    self._process_single_file(
                        example_file,
                        "student_example",
                        quality_label=quality.lower(),
                        assignment_number=assignment_number
                    )

    def _process_folder(self, folder_path: Path, category: str):
        """Process all files in a folder with progress bar"""
        if not folder_path.exists():
            st.warning(f"Folder {folder_path} does not exist")
            return
        
        files = list(folder_path.glob("*"))
        for file_path in tqdm(files, desc=f"Processing {category} files"):
            self._process_single_file(file_path, category)                
    
    def _process_single_file(self, file_path: Path, doc_type: str,
                           quality_label: str = None,
                           assignment_number: str = None):
        """Process a single file with error handling"""
        try:
            content = None
            if file_path.suffix.lower() == '.pdf':
                content = self.processor.read_pdf(str(file_path))
            elif file_path.suffix.lower() == '.docx':
                content = self.processor.read_docx(str(file_path))
            elif file_path.suffix.lower() in ['.pptx', '.ppt']:
                content = self.processor.read_pptx(str(file_path))
            
            if content and content.strip():  # Only process if we have content
                chunks = self.processor.chunk_text(content)
                if chunks:  # Only add if we have valid chunks
                    doc_id = self._generate_doc_id(file_path, doc_type, 
                                                assignment_number, quality_label)
                    self._add_to_index(doc_id, chunks, file_path, doc_type,
                                    quality_label, assignment_number)
            
        except Exception as e:
            self.processor.error_logs.append(f"Error processing {file_path}: {str(e)}")
    
    def _generate_doc_id(self, file_path: Path, doc_type: str,
                        assignment_number: str = None,
                        quality_label: str = None) -> str:
        """Generate unique document ID"""
        doc_id = f"{doc_type}_{file_path.stem}"
        if assignment_number:
            doc_id = f"{assignment_number}_{doc_id}"
        if quality_label:
            doc_id = f"{doc_id}_{quality_label}"
        return doc_id
    
    def _add_to_index(self, doc_id: str, chunks: List[str], 
                     file_path: Path, doc_type: str,
                     quality_label: str = None,
                     assignment_number: str = None):
        """Add document to searchable index"""
        try:
            embeddings = self.embed_model.encode(chunks)
            
            self.document_embeddings[doc_id] = embeddings
            self.document_contents[doc_id] = chunks
            self.document_metadata[doc_id] = {
                'path': str(file_path),
                'type': doc_type,
                'quality_label': quality_label,
                'assignment_number': assignment_number
            }
        except Exception as e:
            self.processor.error_logs.append(f"Error embedding {doc_id}: {str(e)}")
    
    def retrieve_relevant_context(self, 
                                query: str, 
                                top_k: int = 3,
                                include_examples: bool = True) -> List[str]:
        """Retrieve most relevant content for a given query"""
        query_embedding = self.embed_model.encode(query)
        
        all_similarities = []
        for doc_id, embeddings in self.document_embeddings.items():
            # Skip student examples if not requested
            if not include_examples and self.document_metadata[doc_id]['type'] == 'student_example':
                continue
                
            similarities = cosine_similarity([query_embedding], embeddings)[0]
            for idx, sim in enumerate(similarities):
                all_similarities.append((sim, doc_id, idx))
        
        # Sort by similarity and get top_k
        all_similarities.sort(reverse=True)
        relevant_content = []
        
        for sim, doc_id, idx in all_similarities[:top_k]:
            content = self.document_contents[doc_id][idx]
            metadata = self.document_metadata[doc_id]
            source = f"\nSource: {metadata['type']}"
            if metadata['assignment_number']:
                source += f" ({metadata['assignment_number']})"
            if metadata['quality_label']:
                source += f" - {metadata['quality_label']} example"
            relevant_content.append(content + source)
            
        return relevant_content
    
    def generate_response(self, 
                         student_query: str,
                         student_work: str = None,
                         assignment_context: str = None) -> str:
        """Generate tutoring response with error handling"""
        try:
            context = self.retrieve_relevant_context(student_query)
            
            prompt = f"""You are an AI tutor for the course 'Society and the Engineer'.
            
Context from course materials:
{' '.join(context)}"""

            if assignment_context and assignment_context != "None":
                prompt += f"\n\nCurrent assignment: {assignment_context}"

            prompt += f"\n\nStudent query: {student_query}"

            if student_work:
                prompt += f"\n\nStudent work submitted: {student_work}"

            prompt += "\n\nProvide a helpful tutoring response that:"
            prompt += "\n- Addresses the student's specific question"
            prompt += "\n- References relevant course materials"
            prompt += "\n- Provides constructive feedback if student work was submitted"
            prompt += "\n- Encourages critical thinking"
            
            response = self.llm.generate_content(prompt)
            return response.text
            
        except Exception as e:
            error_msg = f"Error generating response: {str(e)}"
            print(error_msg)
            return "I apologize, but I encountered an error generating the response. Please try again or rephrase your question."

def create_tutor_ui():
    st.set_page_config(layout="wide", page_title="ENG3004 AI Tutor")
    st.title("Society and the Engineer - AI Tutor")
    
    if 'tutor' not in st.session_state:
        with st.spinner("Initializing tutor..."):
            try:
                st.session_state.tutor = RAGTutor()
                st.success("Tutor initialized successfully!")
            except Exception as e:
                st.error(f"Error initializing tutor: {str(e)}")
                return
    
    # Assignment context selector
    assignment_options = ["None"] + [
        f"Assignment {i}" for i in range(1, 3) 
        if (Path(COURSE_FOLDER) / f"Assignment{i}").exists()
    ]
    assignment_context = st.selectbox(
        "Select assignment context (if applicable):",
        assignment_options
    )
    
    # Input areas
    student_query = st.text_area("What's your question?")
    student_work = st.text_area("(Optional) Paste your work here for feedback", height=200)
    
    if st.button("Get Help"):
        if student_query:
            with st.spinner("Generating response..."):
                response = st.session_state.tutor.generate_response(
                    student_query, 
                    student_work if student_work else None,
                    assignment_context
                )
                st.write(response)
        else:
            st.warning("Please enter a question")

if __name__ == "__main__":
    create_tutor_ui()